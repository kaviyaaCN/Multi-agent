"""
agents/ppt_agent.py
===================
PPT Generator Agent
===================
Converts the project documentation into a structured PowerPoint
presentation with title, content, and conclusion slides.
Outputs both structured slide data (for display) and a real .pptx file.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from utils.llm_client import get_llm_client
from utils.helpers import extract_json_from_text, truncate_text
from utils.logger import logger


@dataclass
class Slide:
    """Represents a single presentation slide."""
    slide_number: int
    title: str
    bullets: list[str] = field(default_factory=list)
    speaker_notes: str = ""
    slide_type: str = "content"  # "title" | "content" | "section" | "conclusion"


@dataclass
class Presentation:
    """Full presentation container."""
    project_title: str
    slides: list[Slide] = field(default_factory=list)
    theme: str = "professional"

    def summary(self) -> str:
        """Return a plain text outline of the presentation."""
        lines = [f"Presentation: {self.project_title}", f"Total slides: {len(self.slides)}\n"]
        for slide in self.slides:
            lines.append(f"Slide {slide.slide_number}: {slide.title}")
            for b in slide.bullets:
                lines.append(f"  • {b}")
            lines.append("")
        return "\n".join(lines)


_SYSTEM_PROMPT = """You are an expert presentation designer with academic communication skills.
Create clear, concise, and visually descriptive slide content.
Each slide should have a focused title and 4-6 bullet points that are:
- Short (max 10 words each)
- Action-oriented
- Factual and technical
Always respond with ONLY valid JSON — no prose outside the JSON."""

_SLIDES_TEMPLATE = """Convert this academic project documentation into a PowerPoint presentation.

Project Title: {title}
Domain: {domain}

Documentation sections provided:
--- ABSTRACT ---
{abstract}

--- INTRODUCTION ---
{introduction}

--- METHODOLOGY ---
{methodology}

--- RESULTS ---
{results}

--- CONCLUSION ---
{conclusion}

Generate exactly 10 slides in this JSON array format:
[
  {{
    "slide_number": 1,
    "title": "Slide title",
    "slide_type": "title|content|section|conclusion",
    "bullets": ["Point 1", "Point 2", "Point 3", "Point 4"],
    "speaker_notes": "Presenter talking points for this slide"
  }},
  ...
]

Required slide order:
1. Title slide (project name, author placeholder, institution)
2. Agenda / Table of Contents  
3. Introduction & Problem Statement
4. Literature Review highlights
5. Proposed System Architecture
6. Methodology & Implementation
7. Technologies Used
8. Results & Performance Metrics
9. Conclusion & Contributions
10. Future Work & References

Return ONLY the JSON array."""


class PPTGeneratorAgent:
    """
    Agent that generates structured presentation slides from documentation.

    Supports:
    - JSON slide structure for Streamlit display
    - Export to .pptx using python-pptx
    """

    def __init__(self) -> None:
        self._llm = get_llm_client()
        logger.info("PPTGeneratorAgent initialised")

    def generate(
        self,
        topic: str,
        domain: str,
        documentation=None,
    ) -> Presentation:
        """
        Generate a full presentation from project documentation.

        Args:
            topic: Project title.
            domain: Technical domain.
            documentation: ProjectDocumentation instance (optional).

        Returns:
            Presentation instance with all slides.
        """
        logger.info("PPTAgent.generate: topic=%s", topic)

        # Extract section content from documentation if provided
        abstract = methodology = results = introduction = conclusion = ""
        if documentation:
            abstract = truncate_text(documentation.abstract, 600)
            introduction = truncate_text(documentation.introduction, 600)
            methodology = truncate_text(documentation.methodology, 600)
            results = truncate_text(documentation.results_and_discussion, 600)
            conclusion = truncate_text(documentation.conclusion, 400)

        prompt = _SLIDES_TEMPLATE.format(
            title=topic,
            domain=domain,
            abstract=abstract or "No abstract provided.",
            introduction=introduction or "No introduction provided.",
            methodology=methodology or "No methodology provided.",
            results=results or "No results provided.",
            conclusion=conclusion or "No conclusion provided.",
        )

        raw = self._llm.generate(
            prompt=prompt,
            system_instruction=_SYSTEM_PROMPT,
            temperature=0.6,
            max_tokens=3000,
        )

        slides_data = extract_json_from_text(raw)
        if not isinstance(slides_data, list):
            logger.error("PPTAgent: could not parse slides JSON — using fallback")
            slides_data = self._fallback_slides(topic, domain)

        slides = []
        for item in slides_data:
            try:
                slides.append(Slide(
                    slide_number=item.get("slide_number", len(slides) + 1),
                    title=item.get("title", f"Slide {len(slides) + 1}"),
                    bullets=item.get("bullets", []),
                    speaker_notes=item.get("speaker_notes", ""),
                    slide_type=item.get("slide_type", "content"),
                ))
            except Exception as exc:
                logger.warning("Skipping malformed slide: %s", exc)

        presentation = Presentation(project_title=topic, slides=slides)
        logger.info("PPTAgent: generated %d slides", len(slides))
        return presentation

    def export_pptx(
        self,
        presentation: Presentation,
        output_path: Optional[Path] = None,
    ) -> Path:
        """
        Export the Presentation to a real .pptx file using python-pptx.

        Args:
            presentation: Presentation dataclass.
            output_path: Destination path (defaults to data/output/<title>.pptx).

        Returns:
            Path to the created .pptx file.
        """
        try:
            from pptx import Presentation as PptxPresentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError("Install python-pptx: pip install python-pptx")

        prs = PptxPresentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Gamma-like Dark Theme Colors
        BG_COLOR = RGBColor(17, 24, 39)      # Deep Gray/Black
        TITLE_COLOR = RGBColor(181, 113, 255) # Bright Purple (Accent)
        TEXT_COLOR = RGBColor(230, 237, 243)  # Off-White
        FONT_NAME = "Arial"

        for slide in presentation.slides:
            # Apply layout
            layout = prs.slide_layouts[0] if slide.slide_type == "title" else prs.slide_layouts[1]
            sl = prs.slides.add_slide(layout)
            
            # Apply dark background format
            background = sl.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = BG_COLOR

            # Style Title
            title_shape = sl.shapes.title
            if title_shape and title_shape.text_frame:
                title_shape.text = slide.title
                for p in title_shape.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    for run in p.runs:
                        run.font.name = FONT_NAME
                        run.font.color.rgb = TITLE_COLOR
                        run.font.bold = True

            # Style Content / Bullets
            if len(sl.placeholders) > 1 and slide.bullets:
                content_shape = sl.placeholders[1]
                if slide.slide_type == "title":
                    content_shape.text = slide.bullets[0]
                    for p in content_shape.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.LEFT
                        for run in p.runs:
                            run.font.name = FONT_NAME
                            run.font.color.rgb = TEXT_COLOR
                            run.font.size = Pt(20)
                else:
                    tf = content_shape.text_frame
                    tf.clear()
                    for bullet in slide.bullets:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        for run in p.runs:
                            run.font.name = FONT_NAME
                            run.font.color.rgb = TEXT_COLOR
                            run.font.size = Pt(22)

            # Add speaker notes
            if slide.speaker_notes:
                notes_slide = sl.notes_slide
                notes_slide.notes_text_frame.text = slide.speaker_notes

        # Determine output path
        if output_path is None:
            out_dir = Path("data") / "output"
            out_dir.mkdir(parents=True, exist_ok=True)
            safe_name = "".join(c if c.isalnum() else "_" for c in presentation.project_title)
            output_path = out_dir / f"{safe_name}.pptx"

        prs.save(str(output_path))
        logger.info("PPTX exported to %s", output_path)
        return output_path

    # ── Fallback ───────────────────────────────────────────────────────────────

    def _fallback_slides(self, topic: str, domain: str) -> list[dict]:
        """Return a minimal slide structure when LLM parsing fails."""
        return [
            {"slide_number": 1, "title": topic, "slide_type": "title",
             "bullets": [f"Domain: {domain}", "Presented by: [Author Name]"], "speaker_notes": "Title slide"},
            {"slide_number": 2, "title": "Agenda", "slide_type": "content",
             "bullets": ["Introduction", "Methodology", "Results", "Conclusion"], "speaker_notes": ""},
            {"slide_number": 3, "title": "Introduction", "slide_type": "content",
             "bullets": ["Problem statement", "Motivation", "Objectives", "Scope"], "speaker_notes": ""},
            {"slide_number": 4, "title": "Methodology", "slide_type": "content",
             "bullets": ["Data collection", "Preprocessing", "Model architecture", "Training"], "speaker_notes": ""},
            {"slide_number": 5, "title": "Results", "slide_type": "content",
             "bullets": ["Accuracy achieved", "Comparison with baselines", "Key findings"], "speaker_notes": ""},
            {"slide_number": 6, "title": "Conclusion", "slide_type": "conclusion",
             "bullets": ["Summary of contributions", "Limitations", "Future work"], "speaker_notes": ""},
        ]
